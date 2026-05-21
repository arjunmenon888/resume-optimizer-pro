import os
import logging

logger = logging.getLogger(__name__)

# MIME types for each supported format
_PDF_TYPES = {"application/pdf"}
_DOCX_TYPES = {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
_DOC_TYPES = {"application/msword"}
_PPTX_TYPES = {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
_PPT_TYPES = {"application/vnd.ms-powerpoint"}
_TEXT_TYPES = {"text/plain", "text/"}


class ExtractionService:

    def _extract_text(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def _extract_pdf(self, file_path: str) -> str:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n".join(pages).strip()
        if not text:
            raise ValueError("PDF appears to be scanned/image-only. Please use an image file for OCR.")
        return text

    def _extract_docx(self, file_path: str) -> str:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)

    def _extract_pptx(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)

    async def _extract_image(self, file_path: str) -> str:
        from services.ocr_service import ocr_service
        return await ocr_service.extract_from_image(file_path)

    async def extract(self, file_path: str, mime_type: str, filename: str = "") -> str:
        ext = os.path.splitext(filename or file_path)[1].lower()

        if mime_type in _TEXT_TYPES or mime_type.startswith("text/") or ext == ".txt":
            return self._extract_text(file_path)

        if mime_type in _PDF_TYPES or ext == ".pdf":
            return self._extract_pdf(file_path)

        if mime_type in _DOCX_TYPES or ext == ".docx":
            return self._extract_docx(file_path)

        if mime_type in _DOC_TYPES or ext == ".doc":
            raise ValueError("Old .doc format is not supported. Please save the file as .docx and try again.")

        if mime_type in _PPTX_TYPES or ext == ".pptx":
            return self._extract_pptx(file_path)

        if mime_type in _PPT_TYPES or ext == ".ppt":
            raise ValueError("Old .ppt format is not supported. Please save the file as .pptx and try again.")

        if mime_type.startswith("image/") or ext in (".jpg", ".jpeg", ".png", ".webp", ".bmp"):
            return await self._extract_image(file_path)

        raise ValueError(
            f"Unsupported file type '{mime_type or ext}'. "
            "Supported formats: PDF, Word (.docx), PowerPoint (.pptx), Image, Plain text."
        )


extraction_service = ExtractionService()
